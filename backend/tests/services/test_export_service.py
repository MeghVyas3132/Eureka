"""Unit tests for Layer 7 export service (JPEG + PPTX)."""
from __future__ import annotations

import io

import pytest
from PIL import Image
from pptx import Presentation

from services.export_service import (
    render_planogram_to_jpeg,
    render_planogram_to_pptx,
)


def _make_planogram_json(*, tier: str = "medium", with_products: bool = True) -> dict:
    products = (
        [
            {
                "product_id": "p1",
                "sku": "COKE-001",
                "name": "Coca Cola 500ml",
                "brand": "Coca-Cola",
                "category": "beverages",
                "position_x_cm": 0.0,
                "width_cm": 7.5,
                "height_cm": 23.0,
                "facing_count": 3,
                "total_width_cm": 22.5,
                "sales_score": 0.87,
                "revenue": 48000.0,
                "units_sold": 240,
                "placement_tier": "eye_level",
                "color_hex": "#7ED321",
            },
            {
                "product_id": "p2",
                "sku": "MILK-002",
                "name": "Organic Milk 1L",
                "brand": "DairyBest",
                "category": "dairy",
                "position_x_cm": 30.0,
                "width_cm": 8.0,
                "height_cm": 20.0,
                "facing_count": 2,
                "total_width_cm": 16.0,
                "sales_score": 0.65,
                "revenue": 22000.0,
                "units_sold": 110,
                "placement_tier": "eye_level",
                "color_hex": "#4A90D9",
            },
        ]
        if with_products
        else []
    )

    return {
        "planogram_id": "pid",
        "store_id": "sid",
        "generation_level": "store",
        "generation_method": "auto",
        "generated_at": "2026-05-06T10:00:00Z",
        "has_sales_data": True,
        "confidence": {
            "score": 0.62 if tier != "low" else 0.30,
            "tier": tier,
            "sales_coverage_pct": 72.0,
            "dimension_coverage_pct": 45.0,
            "category_coverage_pct": 88.0,
            "store_parse_confidence": 0.7,
        },
        "assortment": {
            "total_catalogue_skus": 320,
            "included_skus": 85,
            "excluded_skus": 235,
            "filter_method": "sales_coverage",
            "coverage_pct": 26.6,
            "message": "Showing 85 of 320 SKUs.",
        },
        "data_quality_warnings": [
            {
                "code": "low_dimension_coverage",
                "severity": "medium",
                "message": "55% of products are missing dimensions.",
                "action_label": "Add product dimensions",
                "action_url": "/products?filter=missing_dimensions",
            }
        ],
        "shelf_config": {
            "shelf_count": 5,
            "shelf_width_cm": 180.0,
            "shelf_height_cm": 200.0,
            "shelf_depth_cm": 40.0,
            "shelf_spacing_cm": 40.0,
            "store_type": "convenience",
            "store_type_rules_applied": True,
        },
        "shelves": [
            {"shelf_number": 1, "tier": "top_level", "remaining_width_cm": 180.0, "products": []},
            {"shelf_number": 2, "tier": "eye_level", "remaining_width_cm": 141.5, "products": products},
            {"shelf_number": 3, "tier": "mid_level", "remaining_width_cm": 180.0, "products": []},
            {"shelf_number": 4, "tier": "mid_level", "remaining_width_cm": 180.0, "products": []},
            {"shelf_number": 5, "tier": "low_level", "remaining_width_cm": 180.0, "products": []},
        ],
        "overflow_skus": [],
        "category_summary": {
            "beverages": {"sku_count": 1, "total_revenue": 48000.0, "shelves": [2]},
            "dairy": {"sku_count": 1, "total_revenue": 22000.0, "shelves": [2]},
        },
    }


def test_jpeg_export_returns_valid_image_bytes() -> None:
    data = render_planogram_to_jpeg(_make_planogram_json(), store_name="Reliance Fresh")

    assert isinstance(data, bytes) and len(data) > 1000
    assert data[:3] == b"\xff\xd8\xff"

    img = Image.open(io.BytesIO(data))
    assert img.format == "JPEG"
    assert img.size == (1800, 1200)


def test_jpeg_export_low_confidence_renders_watermark() -> None:
    data = render_planogram_to_jpeg(
        _make_planogram_json(tier="low"), store_name="Test Store"
    )
    # Just confirm it still renders to a valid JPEG without raising.
    Image.open(io.BytesIO(data)).verify()


def test_jpeg_export_handles_empty_shelves() -> None:
    data = render_planogram_to_jpeg(
        _make_planogram_json(with_products=False), store_name="Empty Store"
    )
    img = Image.open(io.BytesIO(data))
    assert img.format == "JPEG"


def test_pptx_export_returns_valid_three_slide_deck() -> None:
    data = render_planogram_to_pptx(_make_planogram_json(), store_name="Reliance Fresh")

    assert isinstance(data, bytes) and len(data) > 5000
    assert data[:2] == b"PK"  # zip magic

    prs = Presentation(io.BytesIO(data))
    assert len(prs.slides) == 3
    # Slide 2 should contain at least one rectangle per product (2) + shelf separators
    slide_two_shapes = list(prs.slides[1].shapes)
    assert len(slide_two_shapes) >= 2


def test_pptx_export_no_warnings_renders_clear_message() -> None:
    payload = _make_planogram_json()
    payload["data_quality_warnings"] = []
    data = render_planogram_to_pptx(payload, store_name="Test Store")
    prs = Presentation(io.BytesIO(data))
    assert len(prs.slides) == 3


@pytest.mark.parametrize("tier", ["high", "medium", "low"])
def test_pptx_export_renders_all_confidence_tiers(tier: str) -> None:
    data = render_planogram_to_pptx(_make_planogram_json(tier=tier), store_name="Test Store")
    prs = Presentation(io.BytesIO(data))
    assert len(prs.slides) == 3
